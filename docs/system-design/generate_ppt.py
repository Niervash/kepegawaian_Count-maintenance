from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set Content
        tf = slide.placeholders[1].text_frame
        tf.text = content_list[0] if content_list else ""
        
        for i in range(1, len(content_list)):
            p = tf.add_paragraph()
            p.text = content_list[i]
            p.level = 0

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sistem Informasi Kepegawaian (SIMPEG)"
    subtitle.text = "Otomatisasi Monitoring Pangkat & KGB\nMetode Pengembangan: Waterfall"

    # Slide 2: Latar Belakang
    add_slide("Latar Belakang", [
        "Administrasi kepegawaian manual berisiko tinggi terjadi keterlambatan.",
        "Monitoring kenaikan pangkat (4 tahun) dan KGB (2 tahun) sulit dipantau secara real-time.",
        "Pegawai kesulitan melacak progres berkas mereka sendiri."
    ])

    # Slide 3: Solusi Utama
    add_slide("Solusi Sistem", [
        "Dashboard Monitoring otomatis dengan countdown hari.",
        "Sistem Notifikasi H-30, H-14, dan H-7.",
        "Manajemen Role (Admin, Pegawai, Pimpinan) untuk keamanan data.",
        "Digitalisasi dokumen (Upload & Approval)."
    ])

    # Slide 4: Data Flow Diagram (DFD)
    add_slide("Data Flow Diagram (DFD)", [
        "Level 0: Menjelaskan interaksi Pegawai, Admin, dan Pimpinan dengan sistem.",
        "Level 1: Detail proses Login, Tracking, Notifikasi, dan Approval.",
        "(Gunakan gambar dari visual-preview.html untuk visualisasi terbaik)"
    ])

    # Slide 6: Flowchart Sistem
    add_slide("Flowchart Sistem", [
        "Menjelaskan alur kerja dari Login hingga Decision Pimpinan.",
        "Check Role Logic: Memisahkan tampilan berdasarkan hak akses.",
        "Proses Loop: Notifikasi -> Upload -> Verifikasi -> Approval."
    ])

    # Slide 7: Metode Waterfall
    add_slide("Metode Pengembangan: Waterfall", [
        "1. Analysis: Pendefinisian fitur dan kebutuhan user.",
        "2. Design: Pembuatan DFD, Flowchart, dan UI/UX.",
        "3. Coding: Pengembangan menggunakan React & Vite.",
        "4. Testing: Uji coba fitur perhitungan dan notifikasi.",
        "5. Deployment: Implementasi sebagai aplikasi web (SPA)."
    ])

    # Slide 8: Teknologi Terapan
    add_slide("Teknologi yang Digunakan", [
        "Frontend: React 19 & Vite (Performa Tinggi).",
        "Styling: Tailwind CSS v4 (Modern & Responsif).",
        "Routing: TanStack Router (Type-safe Navigation).",
        "Data Storage: LocalStorage (Mock-up Enterprise Ready)."
    ])

    # Slide 9: Kesimpulan
    add_slide("Kesimpulan", [
        "SIMPEG meminimalisir kelalaian administrasi kepegawaian.",
        "Transparansi data bagi pegawai meningkatkan kepuasan kerja.",
        "Sistem siap dikembangkan lebih lanjut dengan integrasi API pihak ketiga."
    ])

    # Slide 10: Terima Kasih
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Terima Kasih"
    subtitle = slide.placeholders[1]
    subtitle.text = "Ada Pertanyaan?\n\nHubungi: [Nama Anda]"

    # Save
    prs.save('SIMPEG_Presentation.pptx')
    print("PPT Berhasil dibuat: SIMPEG_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
